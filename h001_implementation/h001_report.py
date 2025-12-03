# H-001 Automated Insight Engine - Report Generator
# Run this script to generate your PPTX report

import argparse
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

def compute_kpis(df):
    total_impr = int(df['impressions'].sum())
    clicks = int(df['clicks'].sum())
    spend = float(df['spend'].sum())
    ctr = round(clicks / total_impr if total_impr else 0, 4)
    cpc = round(spend / clicks if clicks else 0, 4)
    
    return {
        "impressions": total_impr,
        "clicks": clicks,
        "spend": spend,
        "ctr": ctr,
        "cpc": cpc
    }

def build_ppt(kpis, by_campaign, output_path):
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Automated Insight Engine - Report"
    slide.placeholders[1].text = "Generated Programmatically"

    # KPI slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Metrics Overview"
    tf = slide2.placeholders[1].text_frame
    tf.text = f"Total Impressions: {kpis['impressions']}"
    tf.add_paragraph().text = f"Total Clicks: {kpis['clicks']}"
    tf.add_paragraph().text = f"Total Spend: ${kpis['spend']:.2f}"
    tf.add_paragraph().text = f"CTR: {kpis['ctr']*100:.2f}%"
    tf.add_paragraph().text = f"CPC: ${kpis['cpc']:.2f}"

    # Campaign breakdown
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Campaign Breakdown"
    box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf2 = box.text_frame
    tf2.text = "Campaign | Impressions | Clicks | Spend | CTR | CPC"

    for _, row in by_campaign.iterrows():
        p = tf2.add_paragraph()
        p.text = f"{row['campaign']} | {int(row['impressions'])} | {int(row['clicks'])} | ${row['spend']:.2f} | {(row['clicks']/row['impressions'])*100:.2f}% | ${(row['spend']/row['clicks']):.2f}"

    # Simple narrative slide
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "AI Narrative (Example)"
    slide4.placeholders[1].text = (
        "Campaign A is delivering the strongest performance.\n"
        "CTR indicates healthy engagement. Recommendation:\n"
        "Shift more budget towards top-performing creatives."
    )

    prs.save(output_path)
    print(f"Report saved at: {output_path}")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()

    df = pd.read_csv(args.input)
    kpis = compute_kpis(df)
    by_campaign = df.groupby("campaign", as_index=False).agg({
        "impressions": "sum",
        "clicks": "sum",
        "spend": "sum"
    })

    build_ppt(kpis, by_campaign, args.out)

if __name__ == "__main__":
    main()
