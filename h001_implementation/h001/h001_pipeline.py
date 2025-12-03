#!/usr/bin/env python3
"""
H-001: Automated Insight Engine - simple pipeline
Usage:
    python h001_pipeline.py --data h001/data --out h001/output/h001_report.pptx
"""
import argparse
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import datetime

def ingest_csvs(folder: Path) -> pd.DataFrame:
    """Read all CSVs in folder and return a single DataFrame."""
    csvs = sorted(folder.glob("*.csv"))
    if not csvs:
        raise FileNotFoundError(f"No CSV files found in {folder}")
    dfs = []
    for f in csvs:
        df = pd.read_csv(f, parse_dates=True)
        dfs.append(df)
    df_all = pd.concat(dfs, ignore_index=True, sort=False)
    return df_all

def clean_and_normalize(df: pd.DataFrame) -> pd.DataFrame:
    """Basic cleaning: ensure columns, types, fill missing values."""
    # Standard column names we expect
    expected = ['date','campaign','impressions','clicks','spend']
    # lower-case column names
    df.columns = [c.strip().lower() for c in df.columns]
    # ensure required cols exist
    for c in expected:
        if c not in df.columns:
            # add missing numeric columns as zeros, missing strings as 'unknown'
            if c in ('impressions','clicks','spend'):
                df[c] = 0
            else:
                df[c] = 'unknown'
    # Convert types
    try:
        df['date'] = pd.to_datetime(df['date'], errors='coerce')
    except Exception:
        df['date'] = pd.NaT
    df['impressions'] = pd.to_numeric(df['impressions'], errors='coerce').fillna(0).astype(int)
    df['clicks'] = pd.to_numeric(df['clicks'], errors='coerce').fillna(0).astype(int)
    df['spend'] = pd.to_numeric(df['spend'], errors='coerce').fillna(0.0).astype(float)
    # drop rows with no campaign
    df['campaign'] = df['campaign'].fillna('unknown')
    return df

def compute_kpis(df: pd.DataFrame) -> dict:
    """Compute global KPIs and per-campaign metrics."""
    total_impr = int(df['impressions'].sum())
    total_clicks = int(df['clicks'].sum())
    total_spend = float(df['spend'].sum())
    ctr = (total_clicks / total_impr) if total_impr else 0.0
    cpc = (total_spend / total_clicks) if total_clicks else 0.0
    kpis = {
        'total_impressions': total_impr,
        'total_clicks': total_clicks,
        'total_spend': round(total_spend, 2),
        'ctr': round(ctr, 4),
        'cpc': round(cpc, 4)
    }
    # by campaign
    by_camp = df.groupby('campaign', as_index=False).agg({
        'impressions':'sum','clicks':'sum','spend':'sum'
    })
    by_camp['ctr'] = (by_camp['clicks'] / by_camp['impressions']).fillna(0).round(4)
    by_camp['cpc'] = (by_camp['spend'] / by_camp['clicks']).replace([float('inf')], 0).fillna(0).round(4)
    return kpis, by_camp

def generate_pptx(kpis: dict, by_camp: pd.DataFrame, out_path: Path, title="H-001 Automated Insight Engine Report"):
    """Create a simple PPTX report with the KPIs and table."""
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    # KPI slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Performance Indicators"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = (f"Total Impressions: {kpis['total_impressions']}\n"
               f"Total Clicks: {kpis['total_clicks']}\n"
               f"Total Spend: ${kpis['total_spend']:.2f}\n"
               f"CTR: {kpis['ctr']*100:.2f}%\n"
               f"CPC: ${kpis['cpc']:.2f}")

    # Campaign breakdown slide (text table)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Campaign-level Metrics"
    left = Inches(0.5); top = Inches(1.5); width = Inches(9); height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    header = "Campaign | Impressions | Clicks | Spend | CTR | CPC"
    tf.text = header
    # add each campaign row
    for _, row in by_camp.sort_values('impressions', ascending=False).iterrows():
        line = (f"{row['campaign']} | {int(row['impressions'])} | {int(row['clicks'])} | "
                f"${float(row['spend']):.2f} | {row['ctr']*100:.2f}% | ${row['cpc']:.2f}")
        p = tf.add_paragraph()
        p.text = line

    # Simple narrative slide - placeholder for LLM output
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Automated Narrative (sample)"
    body = slide.shapes.placeholders[1].text_frame
    body.text = ("Summary: The campaigns show the distribution of impressions and clicks across campaigns. "
                 "Consider shifting budget towards campaigns with higher CTR and lower CPC. "
                 "Next steps: run A/B test on top creatives and monitor daily performance.")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True, help="Folder with CSV files")
    parser.add_argument("--out", required=True, help="Output PPTX path")
    args = parser.parse_args()

    data_folder = Path(args.data)
    out_path = Path(args.out)

    df = ingest_csvs(data_folder)
    df = clean_and_normalize(df)
    kpis, by_camp = compute_kpis(df)
    generate_pptx(kpis, by_camp, out_path)
    print(f"Report generated at: {out_path}")

if __name__ == "__main__":
    main()
